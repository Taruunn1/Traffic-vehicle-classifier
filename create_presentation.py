"""
Create PowerPoint Presentation for Traffic Monitoring & Violation Alert System
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor

# Create presentation
prs = Presentation()
prs.slide_width = Inches(10)
prs.slide_height = Inches(7.5)

# Define color scheme
DARK_BG = RGBColor(20, 30, 48)
ACCENT_BLUE = RGBColor(66, 135, 245)
ACCENT_GREEN = RGBColor(52, 211, 153)
ACCENT_ORANGE = RGBColor(249, 115, 22)
WHITE = RGBColor(255, 255, 255)
LIGHT_GRAY = RGBColor(200, 200, 200)

def add_title_slide(prs, title, subtitle):
    """Add a title slide"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = DARK_BG
    
    # Add title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(9), Inches(1.5))
    title_frame = title_box.text_frame
    title_frame.text = title
    title_frame.word_wrap = True
    title_p = title_frame.paragraphs[0]
    title_p.font.size = Pt(54)
    title_p.font.bold = True
    title_p.font.color.rgb = ACCENT_BLUE
    title_p.alignment = PP_ALIGN.CENTER
    
    # Add subtitle
    if subtitle:
        subtitle_box = slide.shapes.add_textbox(Inches(0.5), Inches(4.2), Inches(9), Inches(1.5))
        subtitle_frame = subtitle_box.text_frame
        subtitle_frame.text = subtitle
        subtitle_frame.word_wrap = True
        subtitle_p = subtitle_frame.paragraphs[0]
        subtitle_p.font.size = Pt(28)
        subtitle_p.font.color.rgb = WHITE
        subtitle_p.alignment = PP_ALIGN.CENTER
    
    return slide

def add_content_slide(prs, title, content_list):
    """Add a content slide with bullet points"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = DARK_BG
    
    # Add title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(9), Inches(0.8))
    title_frame = title_box.text_frame
    title_frame.text = title
    title_p = title_frame.paragraphs[0]
    title_p.font.size = Pt(44)
    title_p.font.bold = True
    title_p.font.color.rgb = ACCENT_BLUE
    
    # Add horizontal line
    line = slide.shapes.add_connector(1, Inches(0.5), Inches(1.3), Inches(9.5), Inches(1.3))
    line.line.color.rgb = ACCENT_GREEN
    line.line.width = Pt(2)
    
    # Add content
    content_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.7), Inches(8.4), Inches(5))
    text_frame = content_box.text_frame
    text_frame.word_wrap = True
    
    for i, item in enumerate(content_list):
        if i > 0:
            text_frame.add_paragraph()
        p = text_frame.paragraphs[i]
        p.text = item
        p.font.size = Pt(20)
        p.font.color.rgb = WHITE
        p.space_before = Pt(8)
        p.space_after = Pt(8)
        p.level = 0
    
    return slide

def add_two_column_slide(prs, title, left_title, left_items, right_title, right_items):
    """Add a two-column content slide"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = DARK_BG
    
    # Add title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(9), Inches(0.8))
    title_frame = title_box.text_frame
    title_frame.text = title
    title_p = title_frame.paragraphs[0]
    title_p.font.size = Pt(44)
    title_p.font.bold = True
    title_p.font.color.rgb = ACCENT_BLUE
    
    # Add horizontal line
    line = slide.shapes.add_connector(1, Inches(0.5), Inches(1.3), Inches(9.5), Inches(1.3))
    line.line.color.rgb = ACCENT_GREEN
    line.line.width = Pt(2)
    
    # Left column
    left_title_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.7), Inches(4), Inches(0.4))
    left_title_frame = left_title_box.text_frame
    left_title_frame.text = left_title
    left_title_p = left_title_frame.paragraphs[0]
    left_title_p.font.size = Pt(20)
    left_title_p.font.bold = True
    left_title_p.font.color.rgb = ACCENT_ORANGE
    
    left_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.2), Inches(4.2), Inches(4.8))
    left_frame = left_box.text_frame
    left_frame.word_wrap = True
    for i, item in enumerate(left_items):
        if i > 0:
            left_frame.add_paragraph()
        p = left_frame.paragraphs[i]
        p.text = f"• {item}"
        p.font.size = Pt(16)
        p.font.color.rgb = LIGHT_GRAY
        p.space_before = Pt(6)
        p.space_after = Pt(6)
    
    # Right column
    right_title_box = slide.shapes.add_textbox(Inches(5.2), Inches(1.7), Inches(4), Inches(0.4))
    right_title_frame = right_title_box.text_frame
    right_title_frame.text = right_title
    right_title_p = right_title_frame.paragraphs[0]
    right_title_p.font.size = Pt(20)
    right_title_p.font.bold = True
    right_title_p.font.color.rgb = ACCENT_ORANGE
    
    right_box = slide.shapes.add_textbox(Inches(5.2), Inches(2.2), Inches(4.2), Inches(4.8))
    right_frame = right_box.text_frame
    right_frame.word_wrap = True
    for i, item in enumerate(right_items):
        if i > 0:
            right_frame.add_paragraph()
        p = right_frame.paragraphs[i]
        p.text = f"• {item}"
        p.font.size = Pt(16)
        p.font.color.rgb = LIGHT_GRAY
        p.space_before = Pt(6)
        p.space_after = Pt(6)
    
    return slide

# ===== SLIDE 1: Title Slide =====
add_title_slide(prs, "🚗 Traffic Monitoring", "Vehicle Detection & Violation Alert System")

# ===== SLIDE 2: Project Introduction =====
add_content_slide(prs, "📌 Project Introduction", [
    "🔍 Advanced traffic monitoring system using Deep Learning",
    "🤖 AI-powered vehicle classification (21 vehicle types)",
    "🚨 Real-time traffic violation detection",
    "📊 Comprehensive monitoring and reporting dashboard",
    "⚡ Live camera integration and video analysis",
    "🔐 User-friendly web-based interface"
])

# ===== SLIDE 3: Technology Stack =====
add_two_column_slide(prs, "💻 Technology Stack",
    "Backend Technologies", [
        "Flask - REST API Framework",
        "TensorFlow/Keras - Deep Learning",
        "YOLO - Object Detection",
        "OpenCV - Image Processing",
        "SQLite - Database",
        "Python 3.8+"
    ],
    "Frontend Technologies", [
        "React - User Interface",
        "Axios - HTTP Client",
        "CSS3 - Styling & Animations",
        "Node.js - Runtime",
        "Responsive Design",
        "Dark Theme UI"
    ]
)

# ===== SLIDE 4: Key Tools & Libraries =====
add_content_slide(prs, "🛠️ Tools & Libraries", [
    "📦 Deep Learning: TensorFlow, Keras for vehicle classification",
    "🎯 Object Detection: Ultralytics YOLO for real-time detection",
    "🖼️ Image Processing: OpenCV for video & image analysis",
    "🌐 Web Framework: Flask with REST API for backend",
    "💾 Database: SQLite for models, users, and violation records",
    "🎨 Frontend: React with CSS3 for premium dark theme UI"
])

# ===== SLIDE 5: Project Architecture =====
add_content_slide(prs, "🏗️ Project Architecture", [
    "Backend API (Flask Server) → HTTP REST endpoints on localhost:5000",
    "Database Layer → SQLite for models, violations, user management",
    "ML Models → Pre-trained Keras model & YOLO weights",
    "Frontend UI (React) → Interactive dashboard on localhost:3000",
    "Image/Video Processing → OpenCV pipeline for real-time analysis",
    "File Storage → Uploads folder for images/videos, violations folder for evidence"
])

# ===== SLIDE 6: Core Features =====
add_two_column_slide(prs, "✨ Core Features",
    "Detection & Classification", [
        "Vehicle type classification",
        "Emergency vehicle detection",
        "Real-time webcam streaming",
        "Multi-vehicle detection",
        "Confidence scoring"
    ],
    "Monitoring & Alerts", [
        "Traffic violation detection",
        "Speeding alerts",
        "Red light violation tracking",
        "Wrong lane detection",
        "Violation evidence storage"
    ]
)

# ===== SLIDE 7: How It Works - Workflow =====
add_content_slide(prs, "⚙️ How It Works - Basic Flow", [
    "1️⃣ User uploads image/video or enables camera stream",
    "2️⃣ Backend receives media and preprocesses using OpenCV",
    "3️⃣ YOLO detects vehicles in the frame",
    "4️⃣ Keras model classifies each detected vehicle (21 types)",
    "5️⃣ System analyzes for traffic violations",
    "6️⃣ Results displayed with confidence scores and violation alerts"
])

# ===== SLIDE 8: System Components =====
add_two_column_slide(prs, "🔧 System Components",
    "Backend Components", [
        "app.py - Main Flask API",
        "models_db.py - Model management",
        "violations_db.py - Violation tracking",
        "violation_detector.py - Detection logic",
        "models/ - ML model files",
        "uploads/ - User media storage"
    ],
    "Frontend Components", [
        "ImageUploader - Image selection",
        "CameraStream - Live camera feed",
        "PredictionResult - Results display",
        "ViolationsList - Violation history",
        "ModelSelector - Model switching",
        "ClassesList - Vehicle types"
    ]
)

# ===== SLIDE 9: Data Flow =====
add_content_slide(prs, "📊 Data Processing Pipeline", [
    "Input (Image/Video) → OpenCV Preprocessing",
    "Preprocessing → YOLO Object Detection",
    "Detection → Bounding Boxes with coordinates",
    "Extracted Regions → Keras Classification Model",
    "Classification → Vehicle type + Confidence score",
    "Analysis → Violation detection & Storage → JSON Response to Frontend"
])

# ===== SLIDE 10: Vehicle Classification =====
add_content_slide(prs, "🚗 21 Vehicle Classes", [
    "Sedan, SUV, Truck, Bus, Auto, Bike, Bicycle, Scooter",
    "Car, Ambulance, Fire Truck, Police, Delivery Van",
    "Tempo, Tractor, Rickshaw, Cart, Motorcycle",
    "Heavy Vehicle, Light Vehicle, Others",
    "Each vehicle gets confidence score (0-100%)",
    "Emergency vehicles flagged with special alerts"
])

# ===== SLIDE 11: GUI Features =====
add_content_slide(prs, "🎨 User Interface Features", [
    "🖼️ Image Upload with preview",
    "📹 Real-time camera stream from webcam",
    "📊 Detailed prediction results with confidence bars",
    "📝 Traffic violations history and tracking",
    "🔄 Model selection and switching",
    "🌙 Premium dark theme for comfortable viewing"
])

# ===== SLIDE 12: Setup & Deployment =====
add_two_column_slide(prs, "🚀 Quick Setup",
    "Backend Setup", [
        "1. cd backend",
        "2. pip install -r requirements.txt",
        "3. python app.py",
        "Server runs on port 5000"
    ],
    "Frontend Setup", [
        "1. cd frontend",
        "2. npm install",
        "3. npm start",
        "UI runs on port 3000"
    ]
)

# ===== SLIDE 13: Key Advantages =====
add_content_slide(prs, "💡 Key Advantages", [
    "✅ Real-time processing with GPU/CPU support",
    "✅ Accurate vehicle classification (21 types)",
    "✅ Automated violation detection and recording",
    "✅ User-friendly web interface accessible from browser",
    "✅ Scalable architecture with database integration",
    "✅ Evidence storage for violation verification"
])

# ===== SLIDE 14: Conclusion =====
add_title_slide(prs, "🎯 Thank You!", "Traffic Monitoring & Violation Alert System\nPowered by AI & Deep Learning")

# Save presentation
output_path = r"c:\Users\Admin\OneDrive\Desktop\tarun\Project Traffic Management\Work_10_03\Traffic_Management_Presentation.pptx"
prs.save(output_path)
print(f"✅ Presentation created successfully!")
print(f"📄 Saved to: {output_path}")
print(f"📊 Total slides: {len(prs.slides)}")
